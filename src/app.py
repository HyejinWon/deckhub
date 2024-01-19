# -*- coding: utf-8 -*-

import gradio as gr
from pptx import Presentation
import io
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import openai
import json
import re
import os

openai.api_key = os.environ.get("OPENAI_API_KEY")


def get_filename(file_obj):
    return get_text(file_obj.name)

def get_text(ppt_file):

  msg_fb = os.environ.get("PROMPT")


  result = []

  presentation = Presentation(ppt_file)

  if len(presentation.slides) > 5:
    raise Exception("Max slide length is 5. Please follow the guideline.")

  try:

    for slide_idx, slide in enumerate(presentation.slides): # 슬라이드 마다 읽기
      string = ''

      # get all text in presentation
      # It does not distinguish between the title and the content
      for shape in slide.shapes: # 하나의 슬라이드에 들어간 text 찾기
        if hasattr(shape, "text"):
          string += shape.text + "\n"

      # distinguish title and content
      '''
      # title
      title = slide.shapes.title
      if title is None:
        string += 'title is None'
      else:
        string += title.text
      string += '\n'

      # content
      for idx, content in enumerate(slide.shapes.placeholders):
        # content = slide.shapes.placeholders[1]
        if idx == 0:
          continue
        if content is None:
          string += 'content is None'
        else:
          string += content.text
      '''

      # table
      for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
          table_data = []

          row_count = len(shape.table.rows)
          col_count = len(shape.table.columns)
          for _r in range(0, row_count):
            row = []
            for _c in range(0, col_count):
              cell = shape.table.cell(_r, _c)
              # row 별 데이터를 array로 저장
              row.append(cell.text)
            # row 데이터를 전체 데이터 저장 array에 저장
            table_data.append(row)

          # 필요에 따라서는 pandas의 dataframe 등을 이용해서 데이터 저장
          df_temp = pd.DataFrame(columns=table_data[0], data=table_data[1:])
          string += str(table_data)

      # add delimiter
      string += '\n---------------\n'
      result.append(string)

  except Exception as e:
    result = f"Error: {str(e)}"

  # make script based on slide title & contents
  msg_fb = msg_fb + "- Input:\n" + ' '.join(result)

  response = openai.ChatCompletion.create(
    model="gpt-3.5-turbo",
    messages=[
          {"role": "system", "content": msg_fb},
      ],
      temperature=0.3,
      frequency_penalty=0,
      presence_penalty=0,
      top_p=0.2
  )
  response = response.choices[0].message.content

  done_presentation = preprocessing_script(response, presentation)

  # save pptx file
  done_presentation.save('edited_ppt_file.pptx')

  # return ' '.join(result), 'edited_ppt_file.pptx'
  return response, 'edited_ppt_file.pptx'

def preprocessing_script(response, presentation):
  # split chatgpt's response
  pattern = r"- Slide (\d+) Scripts:"
  segments = re.split(pattern, response)

  # print(response)

  for i in range(1, len(segments), 2): # 몇번째 슬라이드 인지도 나와서 2칸씩 점프
    slide_num = int(segments[i])
    script = segments[i+1].strip()

    # print(script)

    slide = presentation.slides[slide_num - 1]

    # add sclide note each of them
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = script

  return presentation


# with gr.Blocks() as demo:
#   upload_button = gr.UploadButton("Click to Upload a File", file_types=[".pptx"], file_count="multiple")
#   upload_button.upload(get_filename, upload_button, "text")

  # output = gr.Textbox(label='output_box')


demo = gr.Interface(get_filename, "file", outputs=["text", "file"], examples=[["test_final.pptx"]])

demo.launch()
